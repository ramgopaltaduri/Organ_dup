from pptx import Presentation

# load existing presentation to read title slide content
orig = Presentation('final_ppt.pptx')
# create a new presentation from scratch
prs = Presentation()

# copy title slide text from original first slide
orig_title = orig.slides[0]
title_text = ''
subtitle_text = ''
if orig_title.shapes.title:
    title_text = orig_title.shapes.title.text
# try to capture subtitle if exists
for shp in orig_title.shapes:
    if shp != orig_title.shapes.title and hasattr(shp, 'text') and shp.text.strip():
        subtitle_text = shp.text
        break

# add title slide to new prs using layout 0
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = title_text
if subtitle_text and len(slide.shapes) > 1:
    slide.shapes[1].text = subtitle_text

# helper for adding content slides

def add_slide(title, body_lines, layout=1):
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    slide.shapes.title.text = title
    if body_lines:
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                body.text = line
            else:
                p = body.add_paragraph()
                p.text = line
    return slide

# Add new slides per requirements
add_slide('Contents', [
    '1. Abstract',
    '2. Methodology',
    '3. Modules & Gantt Chart',
    '4. UML & Data Flow & Proposed System',
    '5. Equations/Design & Software',
    '6. Algorithms & Techniques',
    '7. Expected Outcomes',
    '8. Plan of Action',
])

add_slide('Abstract', [
    'Organ matching prediction is designed to assist transplant teams by estimating the compatibility and likely outcome of a donor-recipient pair. ',
    'Our models are trained on expansive datasets collected from UNOS and regional registries, enabling them to learn patterns across age groups, blood types, and medical histories.',
    'The system currently supports heart, lung, kidney, and liver predictions with AUC values ranging from 0.62 to 0.84 — demonstrating reliable discrimination between successful and unsuccessful matches.',
    'By relying on data-driven ensembles instead of rigid point-based scores, the platform adapts naturally to new data and reduces human bias.',
])

add_slide('Methodology', [
    'Data collection: obtain donor and recipient records from UNOS and several regional transplant registries, then merge into a unified schema.',
    'Data cleaning: handle missing values, standardize categorical fields (e.g. blood type), and remove obvious outliers.',
    'Feature engineering: create new variables such as age difference, HLA mismatch score, and compute donor viability indicators.',
    'Model training: fit Random Forest, Gradient Boosting, SVM and logistic regression models using stratified cross‑validation for each organ type.',
    'Evaluation: measure AUC, precision, recall, and confusion matrices; perform calibration checks to ensure probabilities are well-behaved.',
    'Deployment: wrap the best-performing ensemble in a REST API and configure a real-time alert subsystem for urgent matches.',
])

add_slide('Modules - Overview', [
    'The system is divided into five major modules, each responsible for a stage of the workflow.',
    'Together they ensure data flows smoothly from raw registry files to actionable predictions for clinicians.',
])
# expand module slides
data_module = add_slide('Module: Data Ingestion', [
    'Responsible for connecting to APIs or reading CSV exports from registries.',
    'Validates incoming records and stores them in a central database for further processing.',
])
preproc_module = add_slide('Module: Preprocessing', [
    'Cleans and normalizes the raw data: fills missing values, encodes categorical features, and standardizes numerical measurements.',
    'Implements feature engineering logic that derives compatibility metrics needed by the models.',
])
model_module = add_slide('Module: Modeling', [
    'Trains multiple machine learning models on historical transplant data, tuning hyperparameters through grid search.',
    'Combines individual learners into an ensemble to improve robustness and accuracy.',
])
interface_module = add_slide('Module: Interface', [
    'Provides a user dashboard built with Streamlit, allowing clinicians to input donor/recipient details and view prediction results.',
    'Also exposes a REST API for integration with hospital information systems.',
])
deploy_module = add_slide('Module: Deployment', [
    'Handles packaging the application as a Docker container and deploying it to a cloud service or on-premise server.',
    'Includes monitoring and logging to track model performance over time.',
])
# add slide for gantt placeholder
add_slide('Gantt Chart', [
    'Timeline of major tasks: data collection, preprocessing, modeling, testing, and deployment.',
    '(Insert actual chart here)',
])

add_slide('UML Diagram', [
    'A UML class diagram captures the main entities such as Donor, Recipient, MatchRecord, and ModelResult.',
    'Each class lists key attributes (e.g. bloodType, age, hlaScore) and associations between donor and recipient.',
    '(Insert actual UML graphic here)',
])
add_slide('Data Flow Diagram', [
    'Illustrates the pipeline: data ingestion → preprocessing → model scoring → results delivery.',
    'Highlights checkpoints for validation and alert generation between stages.',
    '(Insert DFD graphic here)',
])
add_slide('Proposed System', [
    'Centralize diverse transplant data sources to build a more representative training set.',
    'Use ensemble learning (Random Forest + Gradient Boosting) to capture complex compatibility relationships.',
    'Deliver real-time predictions through an alert system that notifies transplant coordinators of high‑priority matches.',
    'Design for scalability so new organs or data sources can be added with minimal re‑engineering.',
])

add_slide('Equations & Design - Part 1', [
    'Key evaluation metrics:',
    '• Accuracy = (TP + TN) / (P + N) — overall correctness across both classes.',
    '• Sensitivity (Recall) = TP / (TP + FN) — ability to identify successful matches.',
    '• Specificity = TN / (TN + FP) — ability to identify poor matches.',
    '• Precision = TP / (TP + FP) — proportion of predicted matches that are correct.',
    '• F1‑Score = 2 × (Precision × Recall) / (Precision + Recall) — harmonic mean for balanced evaluation.',
    '• AUC‑ROC = area under receiver operating characteristic curve; metric for discrimination ability across thresholds.',
])
add_slide('Software & Tools - Part 2', [
    'Backend & ML stack:',
    '• Python 3.11 — primary language for data science and backend logic.',
    '• scikit‑learn — machine learning library for model training, validation, ensemble methods.',
    '• pandas & NumPy — data manipulation, cleaning, and numerical computation.',
    '• Flask or FastAPI — lightweight web framework for REST API endpoints.',
    'Frontend & visualization:',
    '• React.js — modern JavaScript framework for the interactive dashboard UI.',
    '• Streamlit — rapid prototyping dashboard for data visualization and user interaction.',
    '• Matplotlib/Seaborn — statistical plotting for model diagnostics and results.',
    'Deployment & storage:',
    '• MySQL or PostgreSQL — relational database for storing cleaned records and predictions.',
    '• Docker — containerization for reproducible deployment across environments.',
    '• Cloud platforms (AWS/Azure) or on‑premise servers for hosting.',
])
add_slide('Algorithms & Techniques', [
    'Primary algorithms: Random Forest and Gradient Boosting for their ability to handle mixed data types and automatically model interactions.',
    'Supplementary models such as SVM and Logistic Regression are used for comparison and ensembling.',
    'Techniques include SMOTE oversampling to correct class imbalance and recursive feature elimination for selecting the most predictive variables.',
    'Model explainability tools (SHAP, feature importances) help clinicians interpret predictions.',
])

add_slide('Expected Outcomes', [
    'Significantly improved matching accuracy compared to traditional scoring systems, resulting in better patient outcomes.',
    'Faster decision-making by automating compatibility assessments and issuing priority alerts.',
    'A flexible platform that can generalize across heart, lung, kidney, and liver transplants, with capability to incorporate new organ types later.',
    'Provision of actionable insights to transplant centers to optimize resource utilization and reduce organ wastage.',
])

add_slide('Plan of Action', [
    'Complete model evaluation and perform extensive hyperparameter tuning using cross-validation.',
    'Build and refine the front-end (Streamlit dashboard) and back-end API services.',
    'Document methodologies, datasets, and results for the final report and publication.',
    'Package and deploy a working prototype; conduct pilot testing with stakeholders to gather usability feedback.',
])

# save updated presentation (different name if previous file is open)
output_name = 'final_ppt_updated_v2.pptx'
prs.save(output_name)
print(f'Updated presentation saved as {output_name}')
