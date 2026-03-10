from pptx import Presentation
prs = Presentation('final_ppt.pptx')
for i, slide in enumerate(prs.slides):
    print(f'Slide {i+1}:')
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            print(shape.text)
    print('-----')
