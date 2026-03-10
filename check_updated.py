from pptx import Presentation
prs=Presentation('final_ppt_updated.pptx')
for i,slide in enumerate(prs.slides):
    print('Slide', i+1)
    for shape in slide.shapes:
        if hasattr(shape,'text') and shape.text.strip():
            print(shape.text)
    print('----')
